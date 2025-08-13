'''
Quote Generator for Waste Robotics
Author: Cody Martins
'''
import streamlit as st
from PIL import Image
import pandas as pd
from docxtpl import DocxTemplate, InlineImage
from docx.shared import Mm
import os
import matplotlib.pyplot as plt
import tempfile
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import textwrap
import re

# --- WR Branding Setup ---
logo = Image.open("logoWasteRobotics(1).png")  # Make sure this file is in the same directory
col1, col2 = st.columns([1, 6])
with col1:
    st.image("logoWasteRobotics(1).png", width=80)
with col2:
    st.markdown("<h1 style='color: white;'>Waste Robotics Quote Generator</h1>", unsafe_allow_html=True)
    st.markdown("<p style='color: #EF3A2D; font-style: italic;'>Smarter Sorting with Robotics</p>", unsafe_allow_html=True)

def save_df_as_image(df, currency="CAD"):
    df = df.copy()
    df["Unit Price"] = df["Unit Price"].map(lambda x: f"{currency} {x:,.0f}")
    df["Subtotal"] = df["Subtotal"].map(lambda x: f"{currency} {x:,.0f}")

    fig, ax = plt.subplots(figsize=(10, len(df) * 0.5 + 1))
    ax.axis("off")

    table = ax.table(
        cellText=df.values,
        colLabels=df.columns,
        cellLoc='left',
        loc='center'
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)

    for i in range(len(df.columns)):
        table[0, i].set_facecolor("#ef3a2d")
        table[0, i].set_text_props(weight="bold", color="white")

    for row_idx in range(1, len(df) + 1):
        color = "#f2f2f2" if row_idx % 2 == 0 else "#ffffff"
        for col_idx in range(len(df.columns)):
            table[row_idx, col_idx].set_facecolor(color)

    fig.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


st.markdown("""
    <style>
        .reportview-container {
            background-color: #0F0F0F;
            color: white;
        }
        .sidebar .sidebar-content {
            background-color: #1A1A1A;
        }
        h1, h2, h3 {
            color: #EF3A2D;
        }
        .stButton>button {
            background-color: #EF3A2D;
            color: white;
            border: none;
            padding: 0.5em 1em;
            font-weight: bold;
        }
        .stButton>button:hover {
            background-color: #c23024;
        }
        footer {
            visibility: hidden;
        }
        .footer {
            position: fixed;
            bottom: 0;
            width: 100%;
            background-color: #1A1A1A;
            color: white;
            text-align: center;
            padding: 5px;
            font-size: 12px;
        }
    </style>
    """, unsafe_allow_html=True)


# Load pricing data
pricing_df = pd.read_csv("pricing.csv")

# Convert it to a dictionary for easy lookup (now using price_cad, always as float)
PRICING = {
    row["item"]: float(str(row["price_cad"]).replace(",", ""))
    for _, row in pricing_df.iterrows()
}

# --- Constants ---
# All prices in CSV are in CAD, so CAD is the base currency
CURRENCY_CONVERSION = {"CAD": 1.0, "USD": 0.74, "EUR": 0.68}


# --- UI ---
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "Proposal Info", 
    "System Config", 
    "Technical Specs", 
    "Shipping & Timeline", 
    "Inclusions & Quote"
])

with tab1:
    st.header("Proposal Information")
    st.progress(20, text="Step 1 of 5")
    quote_date = st.date_input("Quote Date")
    value_proposition = st.text_input("Value Proposition (Main Proposal Title)")
    client_name = st.text_input("Client Name")
    client_company = st.text_input("Client Company Name")
    salesman_name = st.text_input("Salesperson Name")
    site_location = st.text_input("Site Location")
    shipping_method = st.selectbox("Shipping Method", ["Truck", "Boat"], help="Select the shipping method for delivery.")
    if shipping_method == "Truck":
        num_trucks_or_containers = st.number_input("Number of Trucks", min_value=1, value=1, step=1)
    else:
        num_trucks_or_containers = st.number_input("Number of Containers (Boat)", min_value=1, value=1, step=1)
    currency = st.selectbox("Currency", ["USD", "CAD", "EUR"])
    if currency != "CAD":
        st.markdown(
            f"Get the latest {currency}/CAD exchange rate from [xe.com](https://www.xe.com/currencyconverter/)."
        )
        user_rate = st.number_input(
            f"Enter the current {currency}/CAD exchange rate", min_value=0.0001, value=0.74 if currency == "USD" else 0.68, format="%.4f"
        )
        multiplier = user_rate
    else:
        multiplier = 1.0
    application_overview = st.text_area("Brief Summary of the Application")


with tab2:
    st.header("System Configuration")
    st.progress(40, text="Step 2 of 5")
    materials = st.multiselect("Materials to Sort", ["PCBs", "UBCs", "Trash", "Other"])
    try_and_buy = st.checkbox("Include Try & Buy Option?")
    belt_speed = st.text_input("Belt Speed (m/min)")
    pick_rate = st.text_input("Pick Rate (picks/minute)")
    # Robot Arms (type and quantity)
    robot_types_list = ["Fanuc LR-Mate", "FanucLr10iA", "Fanuc Delta DR3", "Fanuc M10", "Fanuc M20", "Fanuc M710"]
    selected_robot_types = st.multiselect("Robot Arm Types", robot_types_list)
    robot_type = {}
    for rtype in selected_robot_types:
        qty = st.number_input(f"Quantity of {rtype}", min_value=0, value=1, key=f"qty_robot_{rtype}")
        if qty > 0:
            robot_type[rtype] = qty

    # Robot Bases (type and quantity)
    base_types = ["LrMate/Lr10ia", "Delta DR3", "M-10, M-20, M-710"]
    selected_bases = st.multiselect("Robot Base Types", base_types)
    robot_bases = {}
    for base in selected_bases:
        qty = st.number_input(f"Quantity of {base}", min_value=0, value=1, key=f"qty_base_{base}")
        if qty > 0:
            robot_bases[base] = qty

    # Grippers (type and quantity)
    gripper_types_list = ["VentuR", "BagR", "BagR CO", "PinchR Lr & M10", "MonstR", "DagR"]
    selected_grippers = st.multiselect("Gripper Types", gripper_types_list)
    gripper_type = {}
    for gtype in selected_grippers:
        qty = st.number_input(f"Quantity of {gtype}", min_value=0, value=1, key=f"qty_gripper_{gtype}")
        if qty > 0:
            gripper_type[gtype] = qty

    # Consistency check
    total_arms = sum(robot_type.values()) if robot_type else 0
    total_bases = sum(robot_bases.values()) if robot_bases else 0
    total_grippers = sum(gripper_type.values()) if gripper_type else 0
    if total_arms != total_bases or total_arms != total_grippers:
        st.warning(f"⚠️ The total number of robot arms ({total_arms}), robot bases ({total_bases}), and grippers ({total_grippers}) should be the same for a valid configuration.")


with tab3:
    st.header("Technical Specs")
    st.progress(60, text="Step 3 of 5")
    max_object_weight = st.number_input("Maximum Object Weight per Robot (kg)", min_value=0.0)
    # Disposition prompt
    disposition = st.selectbox("Disposition", ["FTF", "IL", "N/A", "QCX"])
    # VRS Model prompt
    vrs_model = st.selectbox("VRS Model", ["900", "1200", "1600", "1800"])
    # Vision System (type and quantity)
    vision_types_list = ["DeepVision System", "HyperVision System"]
    selected_vision_types = st.multiselect("Robot Vision System", vision_types_list)
    vision_system = {}
    for vtype in selected_vision_types:
        qty = st.number_input(f"Quantity of {vtype}", min_value=0, value=1, key=f"qty_vision_{vtype}")
        if qty > 0:
            vision_system[vtype] = qty
    input_power_kva = st.number_input("Input Power (kVA)", min_value=0.0)
    avg_consumption_kw = st.number_input("Average Power Consumption (kW)", min_value=0.0)
    air_consumption_lpm = st.number_input("Total Air Consumption (L/min)", min_value=0)

with tab4:
    st.header("Shipping & Timeline")
    st.progress(80, text="Step 4 of 5")
    # shipping_distance = st.text_input("Estimated Shipping Distance (miles or km)")
    # Shipping distance is now replaced by method/count logic
    order_confirmation_project_kickoff = st.text_input("Order Confirmation / Project Kickoff Duration")
    detailed_engineering = st.text_input("Detailed Engineering Duration")
    engineering_review = st.text_input("Engineering Review Duration")
    procurement_fabrication = st.text_input("Procurement and Fabrication Duration")
    fat_shipping = st.text_input("FAT and Shipping Duration")
    retrofit_installation = st.text_input("Retrofit and Installation Duration")
    commissioning_and_SAT = st.text_input("Commissioning and SAT Duration")
with tab5:

    st.header("Inclusions & Final Quote")
    st.progress(100, text="Step 5 of 5")

    colA, colB = st.columns(2)
    with colA:
        safety_fencing = st.checkbox("Include Safety Fencing?")
        conveyor_var_speed_license = st.checkbox("Include Conveyor Variable Speed License?")
        custom_ai_training = st.checkbox("Include Custom AI Training?")
        robot_validator_license = st.checkbox("Include Robot Validator License?")
        greyparrot_monitoring_unit = st.checkbox("Include GreyParrot Monitoring Unit?")
        installation_supervision = st.checkbox("Include Installation Supervision?")
    with colB:
        additional_sorting_recipes = st.checkbox("Include Additional Sorting Recipes?")
        sat_to_cfa = st.checkbox("Include SAT to CFA?")
        engineering_and_documentation = st.checkbox("Include Engineering & Documentation?")
        online_commissioning = st.checkbox("Include Online Commissioning?")
        installation_commissioning_training = st.checkbox("Include Installation, Commissioning & Training?")
        lips2_support = st.checkbox("Include LIPS2 Support?")
        warranty_option = st.selectbox(
            "Warranty Option",
            ["None", "1 Year (Standard)", "Extended"]
        )

    if st.button("Generate Quote"): 
        # --- Input Validation ---
        missing_fields = []

        # Proposal Info
        if not value_proposition.strip():
            missing_fields.append("Value Proposition")
        if not client_name.strip():
            missing_fields.append("Client Name")
        if not client_company.strip():
            missing_fields.append("Client Company Name")
        if not salesman_name.strip():
            missing_fields.append("Salesperson Name")
        if not site_location.strip():
            missing_fields.append("Site Location")
        if not application_overview.strip():
            missing_fields.append("Application Overview")

        # System Config
        if not materials:
            missing_fields.append("Materials to Sort")
        if not belt_speed.strip():
            missing_fields.append("Belt Speed")
        if not pick_rate.strip():
            missing_fields.append("Pick Rate")
        if not robot_type:
            missing_fields.append("Robot Type")
        if not gripper_type:
            missing_fields.append("Gripper Type")

        # Technical Specs
        if max_object_weight == 0.0:
            missing_fields.append("Max Object Weight")
        if robot_bases == 0:
            missing_fields.append("Number of Robot Bases")
        if not vision_system:
            missing_fields.append("Vision System")
        if input_power_kva == 0.0:
            missing_fields.append("Input Power")
        if avg_consumption_kw == 0.0:
            missing_fields.append("Average Power Consumption")
        if air_consumption_lpm == 0:
            missing_fields.append("Air Consumption")


        # Shipping & Timeline
        # Shipping method/count validation
        if not site_location.strip():
            missing_fields.append("Site Location")
        if not shipping_method:
            missing_fields.append("Shipping Method")
        if not num_trucks_or_containers or num_trucks_or_containers < 1:
            missing_fields.append("Number of Trucks/Containers")
        if not order_confirmation_project_kickoff.strip():
            missing_fields.append("Order Confirmation / Project Kickoff Duration")
        if not detailed_engineering.strip():
            missing_fields.append("Detailed Engineering Duration")
        if not engineering_review.strip():
            missing_fields.append("Engineering Review Duration")
        if not procurement_fabrication.strip():
            missing_fields.append("Procurement and Fabrication Duration")
        if not fat_shipping.strip():
            missing_fields.append("FAT and Shipping Duration")
        if not retrofit_installation.strip():
            missing_fields.append("Retrofit and Installation Duration")
        if not commissioning_and_SAT.strip():
            missing_fields.append("Commissioning and SAT Duration")

        if missing_fields:
            st.error("Please fill in all required fields:\n- " + "\n- ".join(missing_fields))
            st.stop()

        # --- SAFE TO EXECUTE BELOW THIS LINE ---

        # Calculate pricing
        def calculate_price_breakdown(inputs):
            # Only include items that are present in the current PRICING dict (from pricing.csv)
            breakdown = []
            if "Conveyor_Variable_Speed_License" in PRICING and inputs.get("conveyor_var_speed_license"):
                breakdown.append({
                    "Component": "Conveyor Variable Speed License",
                    "Description": "Conveyor Variable Speed License",
                    "Unit Price": PRICING.get("Conveyor_Variable_Speed_License", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("Conveyor_Variable_Speed_License", 0)
                })
            if "custom_ai_training" in PRICING and inputs.get("custom_ai_training"):
                breakdown.append({
                    "Component": "Custom AI Training",
                    "Description": "Custom AI Training",
                    "Unit Price": PRICING.get("custom_ai_training", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("custom_ai_training", 0)
                })
            if "robot_validator_license" in PRICING and inputs.get("robot_validator_license"):
                breakdown.append({
                    "Component": "Robot Validator License",
                    "Description": "Robot Validator License",
                    "Unit Price": PRICING.get("robot_validator_license", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("robot_validator_license", 0)
                })
            if "GreyParrot_Monitoring_Unit" in PRICING and inputs.get("greyparrot_monitoring_unit"):
                breakdown.append({
                    "Component": "GreyParrot Monitoring Unit",
                    "Description": "GreyParrot Monitoring Unit",
                    "Unit Price": PRICING.get("GreyParrot_Monitoring_Unit", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("GreyParrot_Monitoring_Unit", 0)
                })
            if "installation_Supervision" in PRICING and inputs.get("installation_supervision"):
                breakdown.append({
                    "Component": "Installation Supervision",
                    "Description": "Installation Supervision",
                    "Unit Price": PRICING.get("installation_Supervision", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("installation_Supervision", 0)
                })
            if "Additional_Sorting_recipes" in PRICING and inputs.get("additional_sorting_recipes"):
                breakdown.append({
                    "Component": "Additional Sorting Recipes",
                    "Description": "Additional Sorting Recipes",
                    "Unit Price": PRICING.get("Additional_Sorting_recipes", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("Additional_Sorting_recipes", 0)
                })
            if "SAT_to_CFA" in PRICING and inputs.get("sat_to_cfa"):
                breakdown.append({
                    "Component": "SAT to CFA",
                    "Description": "SAT to CFA",
                    "Unit Price": PRICING.get("SAT_to_CFA", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("SAT_to_CFA", 0)
                })
            if "Engineering_&_Documentation" in PRICING and inputs.get("engineering_and_documentation"):
                breakdown.append({
                    "Component": "Engineering & Documentation",
                    "Description": "Engineering & Documentation",
                    "Unit Price": PRICING.get("Engineering_&_Documentation", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("Engineering_&_Documentation", 0)
                })
            if "Online_Commisioning" in PRICING and inputs.get("online_commissioning"):
                breakdown.append({
                    "Component": "Online Commissioning",
                    "Description": "Online Commissioning",
                    "Unit Price": PRICING.get("Online_Commisioning", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("Online_Commisioning", 0)
                })
            if "Installation_Commisioning_&_Training" in PRICING and inputs.get("installation_commissioning_training"):
                breakdown.append({
                    "Component": "Installation, Commissioning & Training",
                    "Description": "Installation, Commissioning & Training",
                    "Unit Price": PRICING.get("Installation_Commisioning_&_Training", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("Installation_Commisioning_&_Training", 0)
                })
            if "LIPS2_support" in PRICING and inputs.get("lips2_support"):
                breakdown.append({
                    "Component": "LIPS2 Support",
                    "Description": "LIPS2 Support",
                    "Unit Price": PRICING.get("LIPS2_support", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("LIPS2_support", 0)
                })

            # Key mappings for CSV
            robot_key_map = {
                "Fanuc LR-Mate": "Fanuc_LR-Mate",
                "FanucLr10iA": "FanucLr10iA",
                "Fanuc Delta DR3": "Fanuc_Delta_DR3",
                "Fanuc M10": "Fanuc_M10",
                "Fanuc M20": "Fanuc_M20",
                "Fanuc M710": "Fanuc_M710"
            }
            gripper_key_map = {
                "VentuR": "VentuR",
                "BagR": "BagR",
                "BagR CO": "BagR_CO",
                "PinchR Lr & M10": "PinchR_Lr_&_M10",
                "MonstR": "MonstR",
                "DagR": "DagR"
            }
            vision_key_map = {
                "DeepVision System": "DeepVision_System",
                "HyperVision System": "HyperVision_System"
            }

            # Robot Base key mapping (update as per your CSV keys)
            base_key_map = {
                "LrMate/Lr10ia": "LrMate/Lr10ia",
                "Delta DR3": "Delta_DR3",
                "M-10, M-20, M-710": "M10_M20_M710"
            }

            # Robot Arms (by type and quantity)
            if isinstance(inputs["robot_type"], dict):
                for rtype, qty in inputs["robot_type"].items():
                    price_key = robot_key_map.get(rtype, rtype)
                    price = PRICING.get(price_key, 0)
                    breakdown.append({
                        "Component": "Robot Arm",
                        "Description": rtype,
                        "Unit Price": price,
                        "Qty": qty,
                        "Subtotal": price * qty
                    })
            else:
                price_key = robot_key_map.get(inputs["robot_type"], inputs["robot_type"])
                price = PRICING.get(price_key, 0)
                breakdown.append({
                    "Component": "Robot Arm",
                    "Description": inputs["robot_type"],
                    "Unit Price": price,
                    "Qty": inputs["robot_arms"],
                    "Subtotal": price * inputs["robot_arms"]
                })

            # Robot Bases (by type and quantity)
            if "robot_bases" in inputs and isinstance(inputs["robot_bases"], dict):
                for btype, qty in inputs["robot_bases"].items():
                    price_key = base_key_map.get(btype, btype)
                    price = PRICING.get(price_key, 0)
                    breakdown.append({
                        "Component": "Robot Base",
                        "Description": btype,
                        "Unit Price": price,
                        "Qty": qty,
                        "Subtotal": price * qty
                    })

            # Grippers (by type and quantity)
            if isinstance(inputs["gripper_type"], dict):
                for gtype, qty in inputs["gripper_type"].items():
                    price_key = gripper_key_map.get(gtype, gtype)
                    price = PRICING.get(price_key, 0)
                    breakdown.append({
                        "Component": "Gripper",
                        "Description": gtype,
                        "Unit Price": price,
                        "Qty": qty,
                        "Subtotal": price * qty
                    })
            else:
                price_key = gripper_key_map.get(inputs["gripper_type"], inputs["gripper_type"])
                price = PRICING.get(price_key, 0)
                breakdown.append({
                    "Component": "Gripper",
                    "Description": str(inputs["gripper_type"]),
                    "Unit Price": price,
                    "Qty": 1,
                    "Subtotal": price
                })

            # Conveyor (only if present in pricing)
            if "conveyor" in PRICING and inputs["conveyor_included"] == "Yes":
                breakdown.append({
                    "Component": "Conveyor",
                    "Description": f"{inputs['conveyor_size']} inch belt",
                    "Unit Price": PRICING.get("conveyor", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("conveyor", 0)
                })

            # Vision Systems (by type and quantity)
            if "vision_system" in inputs and isinstance(inputs["vision_system"], dict):
                for vtype, qty in inputs["vision_system"].items():
                    price_key = vision_key_map.get(vtype, vtype)
                    price = PRICING.get(price_key, 0)
                    breakdown.append({
                        "Component": "Vision System",
                        "Description": vtype,
                        "Unit Price": price,
                        "Qty": qty,
                        "Subtotal": price * qty
                    })


            if inputs.get("try_and_buy"):
                breakdown.append({
                    "Component": "Try & Buy Second Arm",
                    "Description": "Deferred Payment",
                    "Unit Price": PRICING.get("try_and_buy_arm", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("try_and_buy_arm", 0)
                })


            # Shipping logic: by truck or by boat (container)
            shipping_method = inputs.get("shipping_method", "Truck")
            num_units = int(inputs.get("num_trucks_or_containers", 1))
            if shipping_method == "Truck":
                unit_price = PRICING.get("shipping_truck", 8250)
                desc = f"{num_units} truck(s) at ${unit_price:,.0f}/truck"
            else:
                unit_price = PRICING.get("shipping_boat_container", 11000)
                desc = f"{num_units} container(s) at ${unit_price:,.0f}/container (boat)"
            shipping_cost = unit_price * num_units
            breakdown.append({
                "Component": "Shipping",
                "Description": desc,
                "Unit Price": unit_price,
                "Qty": num_units,
                "Subtotal": shipping_cost
            })



            if inputs.get("safety_fencing"):
                breakdown.append({
                    "Component": "Safety Fencing",
                    "Description": "Robot safety fencing",
                    "Unit Price": PRICING.get("safety_fencing", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("safety_fencing", 0)
                })

            # Warranty options
            if inputs["warranty_option"] == "1 Year (Standard)":
                breakdown.append({
                    "Component": "Warranty (1 year)",
                    "Description": "Parts + labor coverage (1 year)",
                    "Unit Price": PRICING.get("warranty_1yr", PRICING.get("warranty", 0)),
                    "Qty": 1,
                    "Subtotal": PRICING.get("warranty_1yr", PRICING.get("warranty", 0))
                })
            elif inputs["warranty_option"] == "Extended":
                breakdown.append({
                    "Component": "Warranty (Extended)",
                    "Description": "Parts + labor coverage (Extended)",
                    "Unit Price": PRICING.get("warranty_extended", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("warranty_extended", 0)
                })

            if inputs.get("pe_stamp"):
                breakdown.append({
                    "Component": "PE Stamp",
                    "Description": "Professional engineer review",
                    "Unit Price": PRICING.get("pe_stamp", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("pe_stamp", 0)
                })

            if inputs.get("sat"):
                breakdown.append({
                    "Component": "Site Acceptance Test (SAT)",
                    "Description": "Final performance check",
                    "Unit Price": PRICING.get("sat", 0),
                    "Qty": 1,
                    "Subtotal": PRICING.get("sat", 0)
                })

            return breakdown

        total_arms = sum(robot_type.values()) if robot_type else 0
        inputs = {
            "robot_arms": total_arms,
            "gripper_type": gripper_type,
            "try_and_buy": try_and_buy,
            "robot_type": robot_type,
            "robot_bases": robot_bases,
            "shipping_method": shipping_method,
            "num_trucks_or_containers": num_trucks_or_containers,
            "safety_fencing": safety_fencing,
            "warranty_option": warranty_option,
            "vision_system": vision_system,
            "conveyor_var_speed_license": conveyor_var_speed_license,
            "custom_ai_training": custom_ai_training,
            "robot_validator_license": robot_validator_license,
            "greyparrot_monitoring_unit": greyparrot_monitoring_unit,
            "installation_supervision": installation_supervision,
            "additional_sorting_recipes": additional_sorting_recipes,
            "sat_to_cfa": sat_to_cfa,
            "engineering_and_documentation": engineering_and_documentation,
            "online_commissioning": online_commissioning,
            "installation_commissioning_training": installation_commissioning_training,
            "lips2_support": lips2_support
        }

        df = pd.DataFrame(calculate_price_breakdown(inputs))
        multiplier = float(CURRENCY_CONVERSION.get(currency, 1.0))
        df["Unit Price"] = pd.to_numeric(df["Unit Price"], errors="coerce").fillna(0) * multiplier
        df["Subtotal"] = pd.to_numeric(df["Subtotal"], errors="coerce").fillna(0) * multiplier

        total = df["Subtotal"].sum()
        st.dataframe(df.style.format({"Unit Price": "${:,.0f}", "Subtotal": "${:,.0f}"}))
        st.markdown(f"### **Total Estimated Price: {currency} {total:,.0f}**")

        # --- DOCX Generation ---
        doc = DocxTemplate("template_practice.docx")

        # Robot Arm Images (one per selected type)
        if robot_type:
            robot_arm_images = []
            for rtype in robot_type.keys():
                base_name = rtype.lower().replace(" ", "_").replace("&", "and").replace(",", "").replace("-", "_")
                robot_arm_filename = f"robot_{base_name}.png"
                if not os.path.exists(robot_arm_filename):
                    robot_arm_filename = "robot_default.png"
                robot_arm_images.append(InlineImage(doc, robot_arm_filename, width=Mm(100), height=Mm(80)))
        else:
            robot_arm_images = [InlineImage(doc, "robot_default.png", width=Mm(100), height=Mm(80))]

        # --- Build Configuration ID for Image Lookup ---
        def sanitize(s):
            return (str(s).strip()
                    .lower()
                    .replace(" ", "_")
                    .replace("&", "and")
                    .replace(",", "")
                    .replace("-", "_")
                    .replace("/", ""))   # ✅ removes slashes


        num_arms = inputs["robot_arms"]

        robot_type_str = "_".join([sanitize(rt) for rt in robot_type.keys()]) if isinstance(robot_type, dict) else sanitize(robot_type)
        disposition_str = sanitize(disposition)
        vrs_model_str = sanitize(vrs_model)
        gripper_type_str = "_".join([sanitize(gt) for gt in gripper_type.keys()]) if isinstance(gripper_type, dict) else sanitize(gripper_type)

        config_id = f"{num_arms}arms_{robot_type_str}_{disposition_str}_{vrs_model_str}_{gripper_type_str}"

        # ✅ Force absolute path for Assets
        base_assets_path = os.path.abspath("Assets")
        assets_folder = os.path.join(base_assets_path, config_id)

        st.write(f"🔍 Looking for config folder: {assets_folder}")  # ✅ Debugging

        # --- Validate the folder ---
        if not os.path.isdir(assets_folder):
            st.warning(f"⚠️ Config folder '{config_id}' not found, using default images.")
            assets_folder = base_assets_path  # fallback to root

        # --- ISO Image ---
        iso_path = os.path.join(assets_folder, "iso.png")
        if not os.path.exists(iso_path):
            st.warning(f"⚠️ Missing iso.png for {config_id}, using default.")
            iso_path = "robot_default.png"
        layout_image = InlineImage(doc, iso_path, width=Mm(100), height=Mm(80))

        # --- Top & Front Images ---
        top_path = os.path.join(assets_folder, "top.png")
        if not os.path.exists(top_path):
            st.warning(f"⚠️ Missing top.png for {config_id}, using default.")
            top_path = "robot_default.png"

        front_path = os.path.join(assets_folder, "front.png")
        if not os.path.exists(front_path):
            st.warning(f"⚠️ Missing front.png for {config_id}, using default.")
            front_path = "robot_default.png"

        layout_overview_top = InlineImage(doc, top_path, width=Mm(150), height=Mm(80))
        layout_overview_front = InlineImage(doc, front_path, width=Mm(150), height=Mm(80))


        if gripper_type:
            gripper_images = []
            for gtype in gripper_type.keys():
                base_name = gtype.lower().replace(" ", "_").replace("&", "and").replace(",", "").replace("-", "_")
                gripper_filename = f"gripper_{base_name}.png"
                if not os.path.exists(gripper_filename):
                    gripper_filename = "gripper_default.png"
                gripper_images.append(InlineImage(doc, gripper_filename, width=Mm(100), height=Mm(80)))
        else:
            gripper_images = [InlineImage(doc, "gripper_default.png", width=Mm(100), height=Mm(80))]


        # Create InlineImage for docxtpl using in-memory BytesIO
        price_table_img = InlineImage(doc, save_df_as_image(df, currency=currency), width=Mm(160))

        context = {
            "value_proposition": value_proposition,
            "application_overview": application_overview,
            "client_name": client_name,
            "client_company": client_company,
            "quote_date": quote_date.strftime("%B %d, %Y"),
            "site_location": site_location,
            "robot_type": robot_type,
            "robot_arms": inputs["robot_arms"],
            "robot_bases": sum(robot_bases.values()) if isinstance(robot_bases, dict) else robot_bases,
            "gripper_type": ", ".join(gripper_type),
            "vision_system": ", ".join(vision_system),
            "materials": ", ".join(materials),
            "belt_speed": belt_speed,
            "pick_rate": pick_rate,
            "max_object_weight": max_object_weight,
            "input_power_kva": input_power_kva,
            "avg_consumption_kw": avg_consumption_kw,
            "air_consumption_lpm": air_consumption_lpm,
            "total_price": f"{currency} {total:,.0f}",
            "warranty_option": warranty_option,
            "safety_fencing": safety_fencing,
            "try_and_buy": try_and_buy,
            "layout_image": layout_image,
            "gripper_images": gripper_images,
            "robot_arm_images": robot_arm_images,
            "order_confirmation_project_kickoff": order_confirmation_project_kickoff,
            "detailed_engineering": detailed_engineering,
            "engineering_review": engineering_review,
            "procurement_fabrication": procurement_fabrication,
            "fat_shipping": fat_shipping,
            "retrofit_installation": retrofit_installation,
            "commissioning_and_SAT": commissioning_and_SAT,
            "price_table_img": price_table_img,
            "layout_overview_top": layout_overview_top,
            "layout_overview_front": layout_overview_front,
        }



        import tempfile

        doc.render(context)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            temp_file_path = tmp.name
            doc.save(temp_file_path)

        st.success("✅ Quote generated successfully!")

        with open(temp_file_path, "rb") as f:
            st.download_button(
                label="📄 Download Quote DOCX",
                data=f,
                file_name=f"{client_name}_Quote_{quote_date.strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

        # --- PowerPoint Generation ---
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        BLUE = RGBColor(46, 125, 122)

        # Helper: Dynamically fit text to box height and width
        def fit_text_to_box(frame, text, box_height_in, box_width_in, max_font=15, min_font=8):
            from pptx.util import Pt
            # Estimate characters per line based on box width (roughly 10pt font = 12 chars/inch)
            chars_per_inch = 12
            for font_size in range(max_font, min_font - 1, -1):
                frame.clear()
                p = frame.add_paragraph()
                p.text = text
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.name = "Arial"
                frame.word_wrap = True
                # Estimate chars per line for this font size
                cpi = chars_per_inch * (font_size / 10)
                max_line_len = int(box_width_in * cpi)
                # Estimate wrapped lines
                lines = []
                for line in text.splitlines():
                    lines.extend(textwrap.wrap(line, width=max_line_len) or [""])
                n_lines = len(lines)
                est_text_height_pt = n_lines * font_size * 1.2
                box_height_pt = box_height_in * 72
                if est_text_height_pt < box_height_pt:
                    break
        def add_page_number(slide, page_num, color=RGBColor(120, 120, 120)):
            slide.shapes.add_textbox(
                slide_width - Inches(1.2),
                slide_height - Inches(0.5),
                Inches(1),
                Inches(0.3)
            ).text_frame.text = f"{page_num}"
            tf = slide.shapes[-1].text_frame
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = color
            p.font.name = FONT_NAME
            p.alignment = 2  # Right

        def add_footer_bar(slide, text="Waste Robotics", color=BLUE):
            bar_height = Pt(4)
            bar = slide.shapes.add_shape(
                1,  # msoShapeRectangle
                0,
                slide_height - bar_height,
                slide_width,
                bar_height
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.width = Pt(0)
            bar.line.fill.background()

        def add_watermark(slide, logo_path="logo2.png"):
            if os.path.exists(logo_path):
                slide.shapes.add_picture(
                    logo_path,
                    slide_width - Inches(2.75),
                    slide_height - Inches(0.5),
                    width=Inches(2),
                    height=Inches(0.25)
                ).element.set('style', 'opacity:0.08')  # Note: python-pptx doesn't support opacity directly, but you can pre-make a transparent PNG.


        # Branding colors
        BRAND_RED = RGBColor(239, 58, 45)   # #EF3A2D
        BRAND_DARK = RGBColor(15, 15, 15)   # #0F0F0F

        def add_branding(slide):
            # Set background color
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = BRAND_DARK

            # Add logo (top left)
            slide.shapes.add_picture("logo1.png", Inches(0.2), Inches(0.2), width=Inches(1.5))

        # Always use the blank layout for new slides
        blank_layout = prs.slide_layouts[-1]  # This is usually the blank slide

        left = Inches(1)
        top = Inches(1.2)
        width = Inches(8)
        height = Inches(1.2)

        # --- Title Slide ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        # Place the image so its left edge is at the center of the slide
        bg_img_path = "title_background.png"
        if os.path.exists(bg_img_path):
            from PIL import Image as PILImage
            img = PILImage.open(bg_img_path)
            img_width, img_height = img.size
            slide_px_height = int(slide_height / 9525)
            # Scale image to fit slide height
            scale = slide_px_height / img_height
            new_width = int(img_width * scale)
            new_height = slide_px_height
            # Place so left edge aligns with center of slide
            left = int(slide_width / 2)
            top = 0
            slide.shapes.add_picture(
                bg_img_path,
                left,
                top,
                width=new_width * 9525,
                height=new_height * 9525
            )

        # Define the light blue color and font name
        LIGHT_BLUE = RGBColor(46, 125, 122)  # #2e7d7a
        FONT_NAME = "Arial"

        # Title text on left half (never overlaps image)
        left = Inches(0.25)
        top = Inches(1.2)
        width = Inches(4.5)
        height = Inches(1.2)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = f"Value Proposition: \n{value_proposition}"
        p.font.size = Pt(30)
        p.font.bold = False
        p.font.color.rgb = LIGHT_BLUE
        p.font.name = FONT_NAME

        # Add a thin line above the info box
        line_left = left
        line_top = top + height + Inches(0.5)
        line_width = width
        line_height = Pt(2)
        line_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            line_left,
            line_top,
            line_width,
            line_height
        )
        fill = line_shape.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_BLUE
        line_shape.line.color.rgb = LIGHT_BLUE
        line_shape.line.width = Pt(0)

        # Info box below title (also only left half)
        info_shape = slide.shapes.add_textbox(left, line_top + Inches(0.2), width, Inches(2.25))
        info_frame = info_shape.text_frame
        info_frame.clear()
        p = info_frame.add_paragraph()
        p.text = f"Presented to: {client_name}\nCompany: {client_company}\n\n\n\n\n\nDate: {quote_date.strftime('%B %d, %Y')}"
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_BLUE
        p.font.name = FONT_NAME

        # --- Application Overview Slide ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        # Colors and font
        BLUE = RGBColor(46, 125, 122)  # #2e7d7a
        WHITE = RGBColor(255, 255, 255)
        FONT_NAME = "Arial"

        # Heading: Application Overview
        heading_left = Inches(0.7)
        heading_top = Inches(1.0)
        heading_width = Inches(7)
        heading_height = Inches(0.8)
        heading_shape = slide.shapes.add_textbox(heading_left, heading_top, heading_width, heading_height)
        heading_frame = heading_shape.text_frame
        heading_frame.clear()
        p = heading_frame.add_paragraph()
        p.text = "Application Overview"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = FONT_NAME

        # Overview text (white)
        overview_left = heading_left
        overview_top = heading_top + heading_height + Inches(0.1)
        overview_width = Inches(7)
        overview_height = Inches(1.2)
        overview_shape = slide.shapes.add_textbox(overview_left, overview_top, overview_width, overview_height)
        overview_frame = overview_shape.text_frame
        overview_frame.clear()
        p = overview_frame.add_paragraph()
        p.text = application_overview
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # "Preliminary layout design (#arm-system)" section
        layout_label_left = heading_left
        layout_label_top = overview_top + overview_height + Inches(0.2)
        layout_label_width = Inches(7)
        layout_label_height = Inches(0.5)
        layout_label_shape = slide.shapes.add_textbox(layout_label_left, layout_label_top, layout_label_width, layout_label_height)
        layout_label_frame = layout_label_shape.text_frame
        layout_label_frame.clear()
        p = layout_label_frame.add_paragraph()
        p.text = f"Preliminary layout design ({inputs['robot_arms']}-arm system)"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = FONT_NAME

        # ISO image centered below the label
        iso_img_top = layout_label_top + layout_label_height + Inches(0.4)
        iso_img_width = Inches(5)
        iso_img_height = Inches(3)
        iso_img_left = heading_left 

        slide.shapes.add_picture(
            iso_path,
            iso_img_left,
            iso_img_top,
            width=iso_img_width,
            height=iso_img_height
        )
        add_page_number(slide, 1)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Layout Images Slide ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        # Title: Layout Overview (#arms-system)
        layout_title = f"Layout Overview ({inputs['robot_arms']}-arm system)"
        title_left = Inches(0.7)
        title_top = Inches(1.0)
        title_width = Inches(7)
        title_height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = layout_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 125, 122)  # #2e7d7a
        p.font.name = "Arial"

        # Arrange images side by side, centered, keeping natural proportions
        from PIL import Image as PILImage

        def get_scaled_size(img_path, max_width_in, max_height_in):
            img = PILImage.open(img_path)
            img_w, img_h = img.size
            dpi = 96  # Assume 96 dpi for conversion
            max_w_px = max_width_in * dpi
            max_h_px = max_height_in * dpi
            scale = min(max_w_px / img_w, max_h_px / img_h, 1.0)
            return img_w * scale / dpi, img_h * scale / dpi  # return in inches

        max_img_width = Inches(4)
        max_img_height = Inches(3)
        spacing = Inches(0.5)

        # Top view image (left)
        top_w_in, top_h_in = get_scaled_size(top_path, 4, 3)
        # Front view image (right)
        front_w_in, front_h_in = get_scaled_size(front_path, 4, 3)

        # Calculate total width for centering
        total_width = Inches(top_w_in) + Inches(front_w_in) + spacing
        img_top = title_top + title_height + Inches(0.3)
        img_left = (slide_width - total_width) // 2

        # Top view image (left)
        slide.shapes.add_picture(
            top_path,
            img_left,
            img_top,
            width=Inches(top_w_in),
            height=Inches(top_h_in)
        )

        # Front view image (right)
        slide.shapes.add_picture(
            front_path,
            img_left + Inches(top_w_in) + spacing,
            img_top,
            width=Inches(front_w_in),
            height=Inches(front_h_in)
        )
        add_page_number(slide, 2)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Robot Arm & Gripper Model Slide ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        BLUE = RGBColor(46, 125, 122)
        FONT_NAME = "Arial"

        # Layout parameters
        margin_left = Inches(0.7)
        margin_right = Inches(5)
        title_top = Inches(1.0)
        title_height = Inches(0.4)
        gap_after_title = Inches(0.4)
        img_top_start = title_top
        img_width = Inches(3.0)
        img_height = Inches(2.2)
        img_spacing = Inches(0.3)

        # --- Robot Arms (left column) ---
        if robot_type:
            for idx, rtype in enumerate(robot_type.keys()):
                # Title for each arm
                arm_title_top = img_top_start + idx * (title_height + img_height + img_spacing)
                arm_title_shape = slide.shapes.add_textbox(
                    margin_left,
                    arm_title_top,
                    img_width,
                    title_height
                )
                arm_title_frame = arm_title_shape.text_frame
                arm_title_frame.clear()
                p = arm_title_frame.add_paragraph()
                p.text = f"Robot Arm Model: {rtype}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = BLUE
                p.font.name = FONT_NAME

                # Image for each arm
                arm_img_top = arm_title_top + title_height + gap_after_title
                base_name = rtype.lower().replace(" ", "_").replace("&", "and").replace(",", "").replace("-", "_")
                robot_arm_filename = f"robot_{base_name}.png"
                if not os.path.exists(robot_arm_filename):
                    robot_arm_filename = "robot_default.png"
                slide.shapes.add_picture(
                    robot_arm_filename,
                    margin_left + Inches(0.15),
                    arm_img_top,
                    width=img_width,
                    height=img_height
                )

        # --- Grippers (right column) ---
        if gripper_type:
            for idx, gtype in enumerate(gripper_type.keys()):
                # Title for each gripper
                gripper_title_top = img_top_start + idx * (title_height + img_height + img_spacing)
                gripper_title_shape = slide.shapes.add_textbox(
                    margin_right - Inches(0.15),
                    gripper_title_top,
                    img_width,
                    title_height
                )
                gripper_title_frame = gripper_title_shape.text_frame
                gripper_title_frame.clear()
                p = gripper_title_frame.add_paragraph()
                p.text = f"Gripper Model: {gtype}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = BLUE
                p.font.name = FONT_NAME

                # Image for each gripper
                gripper_img_top = gripper_title_top + title_height + gap_after_title
                base_name = gtype.lower().replace(" ", "_").replace("&", "and").replace(",", "").replace("-", "_")
                gripper_filename = f"gripper_{base_name}.png"
                if not os.path.exists(gripper_filename):
                    gripper_filename = "gripper_default.png"
                slide.shapes.add_picture(
                    gripper_filename,
                    margin_right,
                    gripper_img_top,
                    width=img_width,
                    height=img_height
                )
        add_page_number(slide, 3)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Vision System Sensor Fusion Slide (Centered Titles/Labels) ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        BLUE = RGBColor(46, 125, 122)
        WHITE = RGBColor(255, 255, 255)
        BLACK = RGBColor(12, 12, 12)
        FONT_NAME = "Arial"

        # --- Centered Main Title ---
        main_title = "ROBOT VISION SYSTEM SENSOR FUSION"
        main_title_width = Inches(8)
        main_title_height = Inches(0.8)
        main_title_left = (slide_width - main_title_width) // 2
        main_title_top = Inches(1)
        title_shape = slide.shapes.add_textbox(main_title_left, main_title_top, main_title_width, main_title_height)
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = main_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = FONT_NAME
        title_frame.paragraphs[0].alignment = 1  # Center

        # Add a thin vertical blue line down the middle below the title
        line_width = Pt(2)
        line_height = Inches(4.5)  # Adjust as needed for your layout
        line_left = (slide_width // 2) - (line_width // 2) - Inches(0.55)
        line_top = main_title_top + main_title_height + Inches(0.1)

        line_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            line_left,
            line_top,
            line_width,
            line_height
        )
        fill = line_shape.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        line_shape.line.color.rgb = BLUE
        line_shape.line.width = Pt(0)

        # --- Left Section: Vision system image ---
        vision_img_path = "vision_system.png"
        vision_img_width = Inches(2.5)
        vision_img_height = Inches(2)
        vision_img_left = Inches(1.4)
        vision_img_top = main_title_top + main_title_height + Inches(1.25)
        slide.shapes.add_picture(
            vision_img_path,
            vision_img_left,
            vision_img_top,
            width=vision_img_width,
            height=vision_img_height
        )

        # Centered Deepvision label below the image
        label_top = vision_img_top + vision_img_height + Inches(0.2)
        label_shape = slide.shapes.add_textbox(
            vision_img_left + Inches(0.6),
            label_top - Inches(0.4),
            vision_img_width,
            Inches(0.4)
        )
        label_frame = label_shape.text_frame
        label_frame.clear()
        p = label_frame.add_paragraph()
        p.text = "Deepvision"
        p.font.size = Pt(16)
        p.font.bold = False
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        label_frame.paragraphs[0].alignment = 1  # Center

        # --- Right Section: Vision comparison image ---
        comparison_img_path = "vision_comparison.png"
        comparison_img_width = Inches(3.75)
        comparison_img_height = Inches(4)
        comparison_img_left = Inches(5.0)
        comparison_img_top = main_title_top + main_title_height + Inches(0.2)
        slide.shapes.add_picture(
            comparison_img_path,
            comparison_img_left,
            comparison_img_top,
            width=comparison_img_width,
            height=comparison_img_height
        )

        # Centered section labels ("Color", "3d", "AI") beneath each third of the comparison image
        section_titles = ["Color", "3d", "AI"]
        section_width_each = comparison_img_width / 3
        section_label_top = comparison_img_top + comparison_img_height + Inches(0.1)
        for i, section in enumerate(section_titles):
            section_left = comparison_img_left + section_width_each * i
            section_shape = slide.shapes.add_textbox(
                section_left  + Inches(0.25),
                section_label_top - Inches(0.4),
                section_width_each,
                Inches(0.3)
            )
            section_frame = section_shape.text_frame
            section_frame.clear()
            p = section_frame.add_paragraph()
            p.text = section
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            section_frame.paragraphs[0].alignment = 1  # Center
        add_page_number(slide, 4)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Inclusions & Exclusions Slide ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)


        BLUE = RGBColor(46, 125, 122)   # #2e7d7a
        RED = RGBColor(239, 58, 45)     # #EF3A2D
        WHITE = RGBColor(255, 255, 255)
        FONT_NAME = "Arial"

        # --- Top black bar ---
        bar_height = Inches(1.25)
        bar_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            0, 0,
            slide_width,
            bar_height
        )
        bar_fill = bar_shape.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = BLACK
        bar_shape.line.width = Pt(0)
        bar_shape.line.fill.background()

        # --- Logo on top left (over the black bar) ---
        logo_path = "logoWasteRobotics(1).png"
        logo_width = Inches(1.5)
        logo_height = Inches(0.8)
        logo_left = Inches(0.2)
        logo_top = Inches(0.1)
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)

        # --- Section backgrounds start below the bar ---
        side_margin = Inches(0.7)
        top_margin = bar_height
        section_width = (slide_width - 2 * side_margin) // 2
        section_height = slide_height - top_margin

        # Slide dimensions
        half_width = slide_width // 2
        slide_height_in = slide_height / 914400  # EMU to inches

        # --- Left: Inclusions ---
        inclusions_left = 0
        inclusions_top = bar_height
        inclusions_width = half_width
        inclusions_height = slide_height - Inches(1.25)

        # Blue background rectangle
        left_bg = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            inclusions_left,
            inclusions_top,
            inclusions_width,
            inclusions_height
        )
        fill = left_bg.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        left_bg.line.width = Pt(0)
        left_bg.line.fill.background()

        # Inclusions label
        label_shape = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.0),
            Inches(3.5),
            Inches(0.6)
        )
        label_frame = label_shape.text_frame
        label_frame.clear()
        p = label_frame.add_paragraph()
        p.text = "Inclusions"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Build inclusions list
        inclusions_list = []
        # #arms of typeofarm
        if robot_type:
            for rtype, qty in robot_type.items():
                inclusions_list.append(f"{qty} x {rtype} robot arm(s)")
        # #robot bases
        if robot_bases:
            for btype, qty in robot_bases.items():
                inclusions_list.append(f"{qty} x {btype} robot base(s)")
        # # of grippers
        if gripper_type:
            for gtype, qty in gripper_type.items():
                inclusions_list.append(f"{qty} x {gtype} gripper(s)")
        # Shipping
        inclusions_list.append(f"Shipping to {site_location}")

        # All selected inclusions from tab 5
        tab5_labels = [
            ("safety_fencing", "Safety Fencing"),
            ("conveyor_var_speed_license", "Conveyor Variable Speed License"),
            ("custom_ai_training", "Custom AI Training"),
            ("robot_validator_license", "Robot Validator License"),
            ("greyparrot_monitoring_unit", "GreyParrot Monitoring Unit"),
            ("installation_supervision", "Installation Supervision"),
            ("additional_sorting_recipes", "Additional Sorting Recipes"),
            ("sat_to_cfa", "SAT to CFA"),
            ("engineering_and_documentation", "Engineering & Documentation"),
            ("online_commissioning", "Online Commissioning"),
            ("installation_commissioning_training", "Installation, Commissioning & Training"),
            ("lips2_support", "LIPS2 Support"),
            ("warranty_option", f"Warranty: {warranty_option}" if warranty_option != "None" else None)
        ]
        for key, label in tab5_labels:
            if key == "warranty_option":
                if warranty_option != "None":
                    inclusions_list.append(label)
            elif locals().get(key):
                inclusions_list.append(label)

        # Inclusions list text box
        inclusions_text = "\n".join(f"• {item}" for item in inclusions_list)
        inclusions_box_left = Inches(0.5)
        inclusions_box_top = Inches(1.7)
        inclusions_box_width = Inches(3.8)
        inclusions_box_height = slide_height - inclusions_box_top - Inches(0.3)  # leave a bottom margin
        inclusions_box = slide.shapes.add_textbox(
            inclusions_box_left,
            inclusions_box_top,
            inclusions_box_width,
            inclusions_box_height
        )
        inclusions_frame = inclusions_box.text_frame
        fit_text_to_box(
            inclusions_frame,
            inclusions_text,
            inclusions_box_height / 914400,  # convert EMU to inches if needed, but Inches() returns EMU
            inclusions_box_width / 914400
        )

        # --- Right: Exclusions ---
        exclusions_left = half_width
        exclusions_top = bar_height
        exclusions_width = half_width
        exclusions_height = slide_height - Inches(1.25);

        # Red background rectangle
        right_bg = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            exclusions_left,
            exclusions_top,
            exclusions_width,
            exclusions_height
        )
        fill = right_bg.fill
        fill.solid()
        fill.fore_color.rgb = RED
        right_bg.line.width = Pt(0)
        right_bg.line.fill.background()
        # Exclusions label
        ex_label_shape = slide.shapes.add_textbox(
            Inches(5.2),
            Inches(1.0),
            Inches(3.5),
            Inches(0.6)
        )
        ex_label_frame = ex_label_shape.text_frame
        ex_label_frame.clear()
        p = ex_label_frame.add_paragraph()
        p.text = "Exclusions"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Build exclusions list (not selected in tab 5)
        exclusions_list = []
        for key, label in tab5_labels:
            if key == "warranty_option":
                if warranty_option == "None":
                    exclusions_list.append("Warranty")
            elif not locals().get(key):
                exclusions_list.append(label)

        # Always include these exclusions
        exclusions_list += [
            "All modifications required on current equipment to integrate the robotic system",
            "Electrical hookup in client’s facility",
            f"Total input power: {input_power_kva}kVA",
            f"Average Power Consumption: {avg_consumption_kw}kW",
            "Internet hookup in client’s facility (up/down 100 Mbits/sec)",
            f"Compressed air hookup in client’s facility (total air consumption: {air_consumption_lpm}L/min)",
            "Taxes, customs and/or duty charges"
        ]

        # Exclusions list text box
        exclusions_text = "\n".join(f"• {item}" for item in exclusions_list)
        exclusions_box_left = Inches(5.2)
        exclusions_box_top = Inches(1.7)
        exclusions_box_width = Inches(3.8)
        exclusions_box_height = slide_height - exclusions_box_top - Inches(0.3)
        exclusions_box = slide.shapes.add_textbox(
            exclusions_box_left,
            exclusions_box_top,
            exclusions_box_width,
            exclusions_box_height
        )
        exclusions_frame = exclusions_box.text_frame
        fit_text_to_box(
            exclusions_frame,
            exclusions_text,
            exclusions_box_height / 914400,
            exclusions_box_width / 914400
        )
        add_page_number(slide, 5)

        # --- System Specifications & Buying Price Slide (Stacked vertically) ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        BLUE = RGBColor(46, 125, 122)
        WHITE = RGBColor(255, 255, 255)
        FONT_NAME = "Arial"

        # Margins and widths
        left_margin = Inches(0.7)
        content_width = slide_width - 2 * left_margin

        # --- System Specifications (top) ---
        specs_top = Inches(1.2)
        specs_label_shape = slide.shapes.add_textbox(
            left_margin,
            specs_top,
            content_width,
            Inches(0.5)
        )
        specs_label_frame = specs_label_shape.text_frame
        specs_label_frame.clear()
        p = specs_label_frame.add_paragraph()
        p.text = "System Specifications"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = FONT_NAME

        specs_content = (
            f"Up to {pick_rate}\n"
            f"Maximum Object Weight Per Robot: {max_object_weight} kg\n"
            f"Robots operating conditions: 5°C to 45°C"
        )
        specs_content_shape = slide.shapes.add_textbox(
            left_margin,
            specs_top + Inches(0.6),
            content_width,
            Inches(1.0)
        )
        specs_content_frame = specs_content_shape.text_frame
        specs_content_frame.clear()
        p = specs_content_frame.add_paragraph()
        p.text = specs_content
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # --- Thin blue line between sections ---
        line_left = left_margin
        line_top = specs_top + Inches(2.7)  # Just below specs section
        line_width = content_width
        line_height = Pt(2)
        line_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            line_left,
            line_top,
            line_width,
            line_height
        )
        fill = line_shape.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        line_shape.line.color.rgb = BLUE
        line_shape.line.width = Pt(0)
        line_shape.line.fill.background()

        # --- Buying Price (below specs) ---
        price_top = line_top + Inches(0.2)
        price_label_shape = slide.shapes.add_textbox(
            left_margin,
            price_top,
            content_width,
            Inches(0.5)
        )
        price_label_frame = price_label_shape.text_frame
        price_label_frame.clear()
        p = price_label_frame.add_paragraph()
        p.text = "Buying Price"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = FONT_NAME

        additional_arm_price = PRICING.get("try_and_buy_arm", 0) * multiplier
        price_content = (
            f"Robotic Sorting System: {currency} {total:,.0f}\n"
            f"Additional Robot Arm: {currency} {additional_arm_price:,.0f}"
        )
        price_content_shape = slide.shapes.add_textbox(
            left_margin,
            price_top + Inches(0.6),
            content_width,
            Inches(1.0)
        )
        price_content_frame = price_content_shape.text_frame
        price_content_frame.clear()
        p = price_content_frame.add_paragraph()
        p.text = price_content
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Disclaimer in small white font
        disclaimer = "* Prices may vary due to exchange rates, inflation, and integration engineering. Valid for 30 days."
        disclaimer_shape = slide.shapes.add_textbox(
            left_margin,
            price_top + Inches(1.7),
            content_width,
            Inches(0.5)
        )
        disclaimer_frame = disclaimer_shape.text_frame
        disclaimer_frame.clear()
        p = disclaimer_frame.add_paragraph()
        p.text = disclaimer
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        add_page_number(slide, 6)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Timeline Slide with Alternating Connectors and Unified Durations ---
        slide = prs.slides.add_slide(blank_layout)
        add_branding(slide)
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        BLUE = RGBColor(46, 125, 122)
        WHITE = RGBColor(255, 255, 255)
        FONT_NAME = "Arial"

        timeline_events = [
            ("Project Kickoff", order_confirmation_project_kickoff),
            ("Detailed Engineering", detailed_engineering),
            ("Engineering Review", engineering_review),
            ("Procurement & Fabrication", procurement_fabrication),
            ("FAT & Shipping", fat_shipping),
            ("Retrofit & Installation", retrofit_installation),
            ("Commissioning \n & SAT", commissioning_and_SAT)
        ]

        timeline_left = Inches(1.0)
        timeline_right = slide_width - Inches(1.0)
        timeline_top = Inches(3.0)
        timeline_height = Pt(3)
        timeline_width = timeline_right - timeline_left

        # Draw timeline line
        line_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            timeline_left,
            timeline_top,
            timeline_width,
            timeline_height
        )
        fill = line_shape.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        line_shape.line.color.rgb = BLUE
        line_shape.line.width = Pt(0)

        n_events = len(timeline_events)
        circle_radius = Pt(14)
        label_box_height = Inches(0.4)
        label_box_width = Inches(1.4)
        connector_length = Inches(0.7)
        connector_width = Pt(2)

        # For unified durations row
        duration_box_height = Inches(0.3)
        duration_row_top = timeline_top + Inches(1.0)
        duration_box_width = Inches(1.0)  # or adjust as needed

        for i, (label, duration) in enumerate(timeline_events):
            x = timeline_left + i * (timeline_width / (n_events - 1))
            y = timeline_top + timeline_height / 2 - circle_radius / 2

            # Draw circle
            circle = slide.shapes.add_shape(
                9,  # msoShapeOval
                x - circle_radius / 2,
                y,
                circle_radius,
                circle_radius
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = BLUE
            circle.line.color.rgb = WHITE
            circle.line.width = Pt(2)

            # Alternate connector direction and label position
            if i % 2 == 0:
                # Upwards connector
                connector_y1 = y
                connector_y2 = y - connector_length
                connector = slide.shapes.add_shape(
                    1,  # msoShapeRectangle
                    x - connector_width / 2,
                    connector_y2,
                    connector_width,
                    connector_y1 - connector_y2
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = BLUE
                connector.line.color.rgb = BLUE
                connector.line.width = Pt(0)
                connector.line.fill.background()

                # Label above connector
                label_shape = slide.shapes.add_textbox(
                    x - label_box_width / 2,
                    connector_y2 - label_box_height - Inches(0.25),
                    label_box_width,
                    label_box_height
                )
                label_frame = label_shape.text_frame
                label_frame.clear()
                p = label_frame.add_paragraph()
                p.text = label
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = FONT_NAME
                label_frame.paragraphs[0].alignment = 1  # Center

                # Duration just below the circle
                duration_shape = slide.shapes.add_textbox(
                    x - duration_box_width / 2,
                    y + circle_radius - Inches(0.25),
                    duration_box_width,
                    duration_box_height
                )
                duration_frame = duration_shape.text_frame
                duration_frame.clear()
                p = duration_frame.add_paragraph()
                p.text = duration
                p.font.size = Pt(12)
                p.font.color.rgb = BLUE
                p.font.name = FONT_NAME
                duration_frame.paragraphs[0].alignment = 1  # Center

            else:
                # Downwards connector
                connector_y1 = y + circle_radius
                connector_y2 = connector_y1 + connector_length
                connector = slide.shapes.add_shape(
                    1,  # msoShapeRectangle
                    x - connector_width / 2,
                    connector_y1,
                    connector_width,
                    connector_y2 - connector_y1
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = BLUE
                connector.line.color.rgb = BLUE
                connector.line.width = Pt(0)
                connector.line.fill.background()

                # Label below connector
                label_shape = slide.shapes.add_textbox(
                    x - label_box_width / 2,
                    connector_y2 - Inches(0.25),
                    label_box_width,
                    label_box_height
                )
                label_frame = label_shape.text_frame
                label_frame.clear()
                p = label_frame.add_paragraph()
                p.text = label
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = FONT_NAME
                label_frame.paragraphs[0].alignment = 1  # Center

                # Duration just above the circle
                duration_shape = slide.shapes.add_textbox(
                    x - duration_box_width / 2,
                    y - duration_box_height - Inches(0.25),
                    duration_box_width,
                    duration_box_height
                )
                duration_frame = duration_shape.text_frame
                duration_frame.clear()
                p = duration_frame.add_paragraph()
                p.text = duration
                p.font.size = Pt(12)
                p.font.color.rgb = BLUE
                p.font.name = FONT_NAME
                duration_frame.paragraphs[0].alignment = 1  # Center

        # --- Thin line beneath the timeline and durations ---
        line_below_top = duration_row_top + duration_box_height + Inches(0.5)
        line_below_left = timeline_left - Inches(0.75)
        line_below_width = timeline_width + Inches(1.5)
        line_below_height = Pt(2)
        line_below = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            line_below_left,
            line_below_top,
            line_below_width,
            line_below_height
        )
        fill = line_below.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        line_below.line.color.rgb = BLUE
        line_below.line.width = Pt(0)
        line_below.line.fill.background()

        # --- Delivery section ---
        # Calculate total weeks (sum numbers in durations)
        def extract_weeks(duration):
            match = re.search(r"(\d+)", str(duration))
            return int(match.group(1)) if match else 0

        total_weeks = sum(extract_weeks(d) for _, d in timeline_events)

        delivery_label_top = line_below_top + Inches(0.3)
        delivery_label_left = timeline_left
        delivery_label_width = timeline_width
        delivery_label_height = Inches(0.4)

        # "Delivery:" label and total weeks
        delivery_shape = slide.shapes.add_textbox(
            delivery_label_left,
            delivery_label_top,
            delivery_label_width,
            delivery_label_height
        )
        delivery_frame = delivery_shape.text_frame
        delivery_frame.clear()
        p = delivery_frame.add_paragraph()
        p.text = f"Delivery: {total_weeks} weeks"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        delivery_frame.paragraphs[0].alignment = 1  # Center

        # Disclaimer
        disclaimer_shape = slide.shapes.add_textbox(
            delivery_label_left,
            delivery_label_top + delivery_label_height,
            delivery_label_width,
            Inches(0.3)
        )
        disclaimer_frame = disclaimer_shape.text_frame
        disclaimer_frame.clear()
        p = disclaimer_frame.add_paragraph()
        p.text = "(to be confirmed at order time)"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        disclaimer_frame.paragraphs[0].alignment = 1  # Center
        add_page_number(slide, 7)
        add_footer_bar(slide)
        add_watermark(slide)

        # --- Download PPTX ---
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            pptx_file_path = tmp.name
            prs.save(pptx_file_path)

        st.success("✅ Quote PowerPoint generated successfully!")

        with open(pptx_file_path, "rb") as f:
            st.download_button(
                label="📊 Download Quote PPTX",
                data=f,
                file_name=f"{client_name}_Quote_{quote_date.strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )


# Footer branding (if needed)
st.markdown("""
    <div class="footer" style='
        position: fixed;
        bottom: 0;
        width: 100%;
        background-color: #1A1A1A;
        color: white;
        text-align: center;
        padding: 5px;
        font-size: 12px;
    '>
        © 2025 Waste Robotics | Internal Tool
    </div>
""", unsafe_allow_html=True)



